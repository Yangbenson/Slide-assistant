import pandas as pd
import collections
import collections.abc
from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Set pandas display options
pd.set_option('display.max_rows', None)
pd.set_option('display.max_columns', None)


def slide_contents(file_path):

    df = pd.DataFrame(columns=["page_num", "titles", "contents"])

    prs = Presentation(file_path)

    for i, slide in enumerate(prs.slides):
        titles = []
        contents = []
        # print(f"Slide {i+1}:")

        for placeholder in slide.placeholders:
            # print(placeholder.placeholder_format.idx)
            # print(placeholder.placeholder_format.type)

            # Check if the placeholder is a title
            if placeholder.placeholder_format.idx == 0:
                titles.append(placeholder.text)
                # print(f"  Title: {placeholder.text}")

            # Check if the placeholder is content
            elif placeholder.placeholder_format.idx == 1:
                contents.append(placeholder.text)
                # print(f"  Content: {placeholder.text}")

        slide = {"page_num": i+1, "titles": titles, "contents": contents}

        df.loc[len(df)] = slide

    return df

slide_contents("PP_samples/5p_sample.pptx")
print(slide_contents("PP_samples/5P_sample.pptx"))
